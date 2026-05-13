"""
Build tamper-detection survey slide deck:
  "Data Tampering Detection: Technologies Across Fields"

Covers how different domains address data tampering, referencing
Merkle-tree-HDF5.tex (S2-D2 research plan).

Run with: python make_tamper_survey_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette ───────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x37, 0x5E)
MID_BLUE   = RGBColor(0x27, 0x5D, 0x9B)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF7)
GOLD       = RGBColor(0xF0, 0xA5, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
RED_CROSS  = RGBColor(0xC0, 0x20, 0x20)
GREEN_CHK  = RGBColor(0x1A, 0x7A, 0x3C)
AMBER      = RGBColor(0xB8, 0x6A, 0x00)
TEAL       = RGBColor(0x1A, 0x6A, 0x6A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def slide_bg(slide, color=DARK_BLUE):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.25, fill=MID_BLUE)
    add_text(slide, title, 0.35, 0.12, 12.5, 0.65,
             size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.72, 12.5, 0.45,
                 size=14, color=LIGHT_BLUE)

def footer(slide, txt="S2-D2 · Data Tampering Detection · Technologies Across Fields"):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill=MID_BLUE)
    add_text(slide, txt, 0.2, 7.17, 12.9, 0.28,
             size=9, color=WHITE, align=PP_ALIGN.CENTER)

def bullet_box(slide, items, l, t, w, h,
               size=15, color=WHITE, indent="  •  "):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = indent + item
        run.font.size = Pt(size)
        run.font.color.rgb = color

# ── SLIDE 1: Title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
add_rect(s, 0, 2.55, 13.33, 0.06, fill=GOLD)
add_text(s, "Data Tampering Detection:",
         0.6, 0.9, 12.0, 0.7, size=30, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, "Technologies Across Fields",
         0.6, 1.6, 12.0, 0.7, size=30, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER)
add_text(s, "From checksums and Merkle trees to ML anomaly detection and post-quantum signatures",
         0.6, 2.75, 12.0, 0.6, size=16, italic=True, color=LIGHT_BLUE,
         align=PP_ALIGN.CENTER)
add_text(s, "Reference: Breitenfeld — Tamper Detection and Data Integrity for Self-Describing Scientific Data (S2-D2)",
         0.6, 3.55, 12.0, 0.5, size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "The HDF Group · S2-D2 Research Plan",
         0.6, 6.5, 12.0, 0.4, size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 2: Agenda ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Overview", "Seven fields — one question: how do they authenticate data?")
footer(s)

agenda = [
    ("Storage & Filesystems",      "ZFS, Btrfs, dm-verity, T10-PI, COW snapshots"),
    ("Software Supply Chain",      "Code signing, Sigstore, Certificate Transparency, package managers"),
    ("Blockchain & Distributed",   "Bitcoin/Ethereum Merkle DAGs, IPFS, LOCKSS"),
    ("Scientific Data Formats",    "HDF5, NetCDF, Zarr, Parquet — and their gaps"),
    ("Machine Learning & AI",      "Data poisoning attacks, ML anomaly detection, QML"),
    ("Cybersecurity & Forensics",  "NTFS journals, IMA/EVM, tamper-tolerant software (TTS)"),
    ("Long-term Archival",         "Post-quantum signatures, hybrid schemes, TSA timestamping"),
]
for i, (field, detail) in enumerate(agenda):
    y = 1.40 + i * 0.82
    add_rect(s, 0.35, y, 4.20, 0.66, fill=MID_BLUE)
    add_text(s, field, 0.50, y + 0.08, 4.0, 0.50,
             size=14, bold=True, color=WHITE)
    add_rect(s, 4.65, y, 8.30, 0.66, fill=RGBColor(0x1E, 0x44, 0x78))
    add_text(s, detail, 4.80, y + 0.12, 8.1, 0.44,
             size=12, color=LIGHT_BLUE)

# ── SLIDE 3: Why it matters ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Why Data Tampering Matters",
           "The threat is real, consequential, and spans every domain that stores data")
footer(s)

examples = [
    (
        "Machine Learning",
        "0.1 % poisoned training samples → > 80 % error-rate increase",
        "Data poisoning attacks corrupt model behavior at minimal adversarial cost.\n"
        "HDF5 files serving as ML training corpora are a direct attack surface.\n"
        "Source: NeurIPS 2020 data-poisoning study",
        RGBColor(0x7A, 0x20, 0x20),
    ),
    (
        "Scientific Archives",
        "Datasets retained 30–100 years across shared filesystems, tape, and cloud",
        "Silent bit-rot, malicious insiders, in-transit attackers, and future quantum\n"
        "adversaries can all modify data. Current HDF5 has no cryptographic auth.\n"
        "Source: §1 — S2-D2 research plan",
        RGBColor(0x1A, 0x50, 0x30),
    ),
    (
        "Software Supply Chain",
        "SolarWinds, XZ Utils: build artifacts tampered before distribution",
        "Unsigned or weakly authenticated binaries enable nation-state attacks.\n"
        "Code signing and transparency logs now mandatory in major ecosystems.\n"
        "Source: §3.6 — hash chains and append-only transparency logs",
        AMBER,
    ),
    (
        "Forensics & Legal",
        "Tampered evidence or sensor logs — chain of custody requires proof of integrity",
        "NTFS journals, digital signatures, and trusted timestamps are used\n"
        "to prove that files have not been modified since collection.\n"
        "Source: tamper-detection notes — forensic journal analysis",
        RGBColor(0x1A, 0x3A, 0x6A),
    ),
]
for i, (field, headline, body, color) in enumerate(examples):
    col = i % 2
    row = i // 2
    x = 0.30 + col * 6.52
    y = 1.38 + row * 2.75
    add_rect(s, x, y, 6.25, 2.60, fill=color)
    add_text(s, field,    x+0.15, y+0.08,  6.0, 0.38, size=16, bold=True, color=GOLD)
    add_text(s, headline, x+0.15, y+0.48,  6.0, 0.36, size=12, bold=True, color=WHITE)
    add_text(s, body,     x+0.15, y+0.90,  6.0, 1.60, size=11, color=LIGHT_BLUE)

# ── SLIDE 4: Technology taxonomy ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_GRAY)
header_bar(s, "Technology Landscape: Detection Approaches",
           "Ordered by strength of guarantee — from error detection to cryptographic binding")
footer(s)

techs = [
    ("Error-detection codes",   "CRC-32, Fletcher, Adler, Reed-Solomon, T10-PI",
     "Non-adversarial only", "✗", "✗", "✗"),
    ("Cryptographic hashes",    "SHA-256, BLAKE3, SHA-3, KangarooTwelve",
     "Adversarial w/ signature", "✓", "✗", "✗"),
    ("MACs / AEAD",             "HMAC-SHA-256, ChaCha20-Poly1305, AES-GCM",
     "Shared-key auth", "✓", "✗", "✗"),
    ("Digital signatures",      "Ed25519, ECDSA, RSA-PSS, ML-DSA, SLH-DSA",
     "Non-repudiable", "✓", "✓", "Classical: ✗  PQ: ✓"),
    ("Merkle trees",            "Git, ZFS, Bitcoin, IPFS, HDF5 (proposed)",
     "Subset proofs", "✓", "✓", "✓"),
    ("Blockchain / logs",       "Bitcoin, Ethereum, Sigstore Rekor, CT logs",
     "Append-only ledger", "✓", "✓", "Classical: ✗"),
    ("ML anomaly detection",    "Outlier scoring, QML, neural detectors",
     "Probabilistic", "~", "—", "—"),
    ("Forensic journals",       "NTFS $UsnJrnl, Linux IMA, $LogFile",
     "Post-hoc attribution", "~", "✗", "—"),
    ("Tamper-tolerant SW (TTS)","Obfuscation, self-checks, graceful degrade",
     "Tolerance, not proof", "~", "—", "—"),
]

col_labels = ["Technology", "Examples", "Guarantee type", "Adv?", "Public?", "PQ?"]
col_x = [0.20, 2.50, 6.10, 8.85, 9.75, 10.80]
col_w = [2.20, 3.50, 2.65, 0.80, 0.95, 1.90]
ROW_H = 0.53

header_y = 1.32
for j, (lbl, x, w) in enumerate(zip(col_labels, col_x, col_w)):
    add_rect(s, x, header_y, w - 0.04, 0.42, fill=MID_BLUE)
    add_text(s, lbl, x+0.04, header_y+0.04, w-0.08, 0.34,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, (name, examples, gtype, adv, pub, pq) in enumerate(techs):
    y = header_y + 0.44 + i * ROW_H
    bg = RGBColor(0xEB, 0xF2, 0xFB) if i % 2 == 0 else WHITE
    add_rect(s, col_x[0], y, sum(col_w) + len(col_w)*0.04 - 0.04, ROW_H - 0.04, fill=bg)
    cells = [name, examples, gtype, adv, pub, pq]
    for j, (cell, x, w) in enumerate(zip(cells, col_x, col_w)):
        c = DARK_BLUE
        if cell == "✓": c = GREEN_CHK
        elif cell == "✗": c = RED_CROSS
        elif cell == "~": c = AMBER
        bold = j == 0
        add_text(s, cell, x+0.04, y+0.06, w-0.08, ROW_H-0.12,
                 size=9, bold=bold, color=c,
                 align=PP_ALIGN.CENTER if j > 1 else PP_ALIGN.LEFT)

add_text(s, "Adv = defeats adversarial tampering  |  Public = verifiable without secret key  |  PQ = post-quantum safe",
         0.20, 7.00, 12.9, 0.28, size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 5: Storage & Filesystems ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Storage & Filesystems",
           "The most mature integrity infrastructure — but guarantees are stripped in transfer")
footer(s)

storage_techs = [
    ("T10-PI / CRC",
     "SCSI Protection Information (T10-DIF): hardware-accelerated end-to-end CRC\n"
     "from compute node to physical drive. Used in enterprise SAN/NVMe.\n"
     "Non-adversarial only — linear codes are invertible by an adversary."),
    ("ZFS block Merkle",
     "Fletcher4 or SHA-256 checksum per block, arranged in a Merkle-style hierarchy.\n"
     "Enables silent-corruption detection and self-healing from redundant copies.\n"
     "Also used internally by Git (content-addressed Merkle DAG) and IPFS."),
    ("Copy-on-Write (COW)",
     "Both ZFS and Btrfs write mutations as new blocks, preserving immutable snapshots.\n"
     "COW enables point-in-time rollback and self-healing — but does NOT authenticate.\n"
     "An adversary with write access produces a valid COW snapshot of tampered data."),
    ("dm-verity / fs-verity",
     "Linux kernel-enforced Merkle hashes bound to read-time policy (Android Verified Boot).\n"
     "fs-verity signs a per-file Merkle tree; reads are verified on-the-fly by the kernel.\n"
     "Vendor-locked: guarantee evaporates when the file is copied off the protected volume."),
    ("DAOS / Lustre",
     "Exascale HPC object stores: DAOS provides native end-to-end AEAD integrity over NVMe.\n"
     "Lustre uses Adler-32 / CRC-32 wire checksums with T10-PI hardware acceleration.\n"
     "HPC-scale but infrastructure-local — no portability across institutions."),
]

for i, (title, body) in enumerate(storage_techs):
    col = i % 2
    row = i // 2
    if i == 4:
        x, w = 0.30, 12.70
    else:
        x = 0.30 + col * 6.55
        w = 6.30
    y = 1.38 + row * 1.78
    add_rect(s, x, y, w, 1.65, fill=MID_BLUE)
    add_text(s, title, x+0.15, y+0.08, w-0.25, 0.36, size=13, bold=True, color=GOLD)
    add_text(s, body,  x+0.15, y+0.46, w-0.25, 1.10, size=10, color=LIGHT_BLUE)

add_rect(s, 0, 6.92, 13.33, 0.23, fill=AMBER)
add_text(s, "Key limitation: integrity is tied to the storage substrate. Copy to S3 / tape / collaborator laptop — guarantee is gone.",
         0.20, 6.93, 12.9, 0.22, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── SLIDE 6: Software Supply Chain ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Software Supply Chain",
           "Code signing and transparency logs now mandatory — lessons applicable to scientific data")
footer(s)

sw_items = [
    ("Code signing certificates",
     "Windows Authenticode, Apple notarization, RPM GPG.\n"
     "Each binary/package is signed with a private key; verifiers check against a trusted CA.\n"
     "Per-artifact signatures — no subset proof needed since the artifact is atomic."),
    ("Package manager Merkle trees",
     "npm, Cargo, PyPI, apt — each registry uses a Merkle tree over all published packages.\n"
     "A single signed root commits to the entire registry; clients verify inclusion proofs.\n"
     "Directly analogous to the HDF5 Merkle-tree proposal at the file level."),
    ("Certificate Transparency (CT)",
     "Append-only Merkle log of all TLS certificates (RFC 6962).\n"
     "Any certificate not in the public log is rejected by browsers.\n"
     "Tamper-evident: altering a past entry breaks every subsequent root hash."),
    ("Sigstore Rekor",
     "Sigstore extends CT to software artifacts: binaries, container images, SBOMs.\n"
     "Signers get a transparency log inclusion proof; verifiers check it without a PKI.\n"
     "Proposed as a freshness layer for HDF5 signed roots (§5.9 of S2-D2 plan)."),
    ("Secure Boot / UEFI",
     "Platform firmware chain-of-trust: each boot stage verifies the next against a signature.\n"
     "TPM 2.0 extends measured boot to OS and application layers.\n"
     "Hardware-rooted attestation — vendor-locked but extremely strong within its trust domain."),
    ("SBOM + VEX",
     "Software Bill of Materials (SBOM) lists all components; VEX documents known vulnerabilities.\n"
     "Integrity of the SBOM itself must be signed — otherwise an adversary can strip vulnerability data.\n"
     "Emerging standard: CycloneDX, SPDX signed with Sigstore."),
]

for i, (title, body) in enumerate(sw_items):
    col = i % 2
    row = i // 2
    x = 0.30 + col * 6.55
    y = 1.38 + row * 1.82
    add_rect(s, x, y, 6.30, 1.68, fill=RGBColor(0x22, 0x45, 0x70))
    add_text(s, title, x+0.15, y+0.08, 6.05, 0.36, size=13, bold=True, color=GOLD)
    add_text(s, body,  x+0.15, y+0.46, 6.05, 1.14, size=10, color=LIGHT_BLUE)

# ── SLIDE 7: Blockchain & Distributed Systems ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Blockchain & Distributed Systems",
           "Merkle trees are the integrity primitive inside every major blockchain — inherited without the consensus layer")
footer(s)

bc_items = [
    ("Bitcoin",
     "Each block header commits to a Merkle tree of all transactions.\n"
     "Light clients verify transaction inclusion with O(log N) proofs.\n"
     "SHA-256d — not post-quantum, but consensus makes forgery computationally infeasible."),
    ("Ethereum",
     "Merkle-Patricia trie over world state (accounts, storage, code).\n"
     "State root in each block header; ZK rollups use KZG/STARK proofs over state transitions.\n"
     "Verkle tree migration underway for shorter witness sizes."),
    ("IPFS / IPLD",
     "Content-addressed Merkle DAG: every object identified by its hash (CID).\n"
     "Retrieving a CID guarantees content integrity — any corruption changes the address.\n"
     "Used for archival of NFT metadata, scientific datasets, and Web3 applications."),
    ("LOCKSS",
     "Lots of Copies Keep Stuff Safe — peer-to-peer digital preservation.\n"
     "Redundant replicas perform mutual integrity audits (polling-based consensus).\n"
     "Used by university libraries to preserve electronic journals and government records."),
    ("Merkle Mountain Ranges (MMR)",
     "Append-friendly alternative to balanced binary Merkle trees.\n"
     "Used in Grin (cryptocurrency) and proposed for Certificate Transparency v2.\n"
     "Better amortized append cost than balanced trees — relevant to sensor-log archives."),
    ("Key insight for scientific data",
     "All blockchain integrity rests on a Merkle tree.\n"
     "The consensus layer (mining, PoS) is needed for distributed coordination.\n"
     "For a closed-organization archive, the Merkle tree can be inherited without consensus."),
]

for i, (title, body) in enumerate(bc_items):
    col = i % 2
    row = i // 2
    x = 0.30 + col * 6.55
    y = 1.38 + row * 1.82
    is_key = "Key insight" in title
    bg = GREEN_CHK if is_key else RGBColor(0x22, 0x45, 0x70)
    add_rect(s, x, y, 6.30, 1.68, fill=bg)
    add_text(s, title, x+0.15, y+0.08, 6.05, 0.36,
             size=13, bold=True, color=GOLD)
    add_text(s, body,  x+0.15, y+0.46, 6.05, 1.14,
             size=10, color=LIGHT_BLUE if not is_key else WHITE)

# ── SLIDE 8: Scientific Data Formats ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_GRAY)
header_bar(s, "Scientific Data Formats",
           "Existing formats protect against accidental corruption — none provide cryptographic authentication")
footer(s)

formats = [
    ("HDF5",         "Jenkins lookup3 on metadata pages; fletcher32 chunk filter",
     "Non-adversarial CRC only. No per-chunk auth, no subset proofs, no signing."),
    ("NetCDF / CDF-5","CDF-5 adds 64-bit variables; external integrity pipelines explored",
     "No sub-file granularity integrity. No adversarial protection."),
    ("Zarr v3",       "Per-chunk CRC-32C or SHA-256 codecs per chunk",
     "Flat hashes — no hierarchical aggregation enabling subset coverage proofs."),
    ("Apache Parquet","Per-page CRC-32 checksums",
     "Non-adversarial. No signing, no subset proofs, no PQ path."),
    ("ASDF",          "Per-block checksums",
     "Flat only — no Merkle tree, no signatures."),
    ("FITS / DICOM",  "Weak or absent integrity primitives",
     "No cryptographic integrity at all. Widely used in astronomy and medical imaging."),
    ("TileDB",        "AES-256-GCM at-rest tile encryption",
     "Relies on cloud store for object integrity. No verifiable-subset capability."),
    ("HDF5 + Merkle\n(proposed S2-D2)","Merkle tree over chunks + hybrid PQ signature",
     "Subset proofs, tamper localization, incremental update, post-quantum. First of its kind."),
]

col_x = [0.22, 2.90, 6.10]
col_w = [2.58, 3.10, 4.95]
header_y = 1.30
col_labels = ["Format", "Current integrity", "Gap"]
for j, (lbl, x, w) in enumerate(zip(col_labels, col_x, col_w)):
    add_rect(s, x, header_y, w - 0.04, 0.38, fill=MID_BLUE)
    add_text(s, lbl, x+0.06, header_y+0.04, w-0.10, 0.30,
             size=11, bold=True, color=WHITE)

ROW_H = 0.60
for i, (fmt, current, gap) in enumerate(formats):
    y = header_y + 0.40 + i * ROW_H
    is_proposed = "proposed" in fmt
    bg = RGBColor(0xC8, 0xE8, 0xC8) if is_proposed else (
         RGBColor(0xEB, 0xF2, 0xFB) if i % 2 == 0 else WHITE)
    add_rect(s, col_x[0], y, sum(col_w) + 0.12, ROW_H - 0.04, fill=bg)
    cells = [fmt, current, gap]
    for j, (cell, x, w) in enumerate(zip(cells, col_x, col_w)):
        c = DARK_BLUE if not is_proposed else RGBColor(0x0A, 0x50, 0x20)
        if j == 2 and not is_proposed:
            c = RED_CROSS
        if j == 2 and is_proposed:
            c = GREEN_CHK
        add_text(s, cell, x+0.06, y+0.06, w-0.10, ROW_H-0.12,
                 size=9, bold=(j == 0), color=c)

# ── SLIDE 9: Machine Learning & AI ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Machine Learning & AI",
           "ML datasets are a high-value attack surface — and ML itself is used for detection")
footer(s)

add_rect(s, 0.35, 1.38, 12.63, 0.88, fill=RGBColor(0x7A, 0x20, 0x20))
add_text(s,
         "Data poisoning threat:  0.1 % of corrupted training samples → > 80 % increase in model error rate",
         0.55, 1.44, 12.2, 0.38, size=15, bold=True, color=WHITE)
add_text(s,
         "When HDF5 files serve as training corpora, cryptographic integrity is a prerequisite for trustworthy ML. (NeurIPS 2020)",
         0.55, 1.80, 12.2, 0.38, size=12, italic=True, color=RGBColor(0xFF, 0xCC, 0xCC))

ml_items = [
    ("Attack: Data Poisoning",
     "Adversary injects mislabeled or corrupted samples into a training dataset.\n"
     "Targeted (backdoor) attacks alter model behavior on specific inputs only.\n"
     "0.1 % corruption is enough for catastrophic model degradation (NeurIPS 2020)."),
    ("Attack: Model Inversion",
     "Adversary queries a trained model to reconstruct training data.\n"
     "Integrity protection of training sets limits the attack surface.\n"
     "Not directly a tampering attack, but motivates data provenance."),
    ("Defense: ML Anomaly Detection",
     "Outlier detectors flag samples that deviate statistically from the training distribution.\n"
     "Reported accuracy: 75–95 % for targeted poisoning attacks.\n"
     "Probabilistic — cannot provide cryptographic binding to a signer."),
    ("Defense: QML Detection",
     "Quantum ML approaches (quantum amplitude estimation, quantum kernel methods)\n"
     "accelerate inner-product computation for nearest-neighbor outlier scoring.\n"
     "Active research; not yet production-grade for large-scale scientific arrays."),
    ("Defense: Cryptographic Auth",
     "Hash the dataset before training; verify before use.\n"
     "Merkle tree enables per-chunk verification without full re-read.\n"
     "Only approach that provides a proof, not just a signal."),
    ("Integration opportunity",
     "ML anomaly detection + Merkle authentication are complementary.\n"
     "Merkle detects ANY byte change; ML detects semantically suspicious patterns.\n"
     "Neither alone is sufficient for high-stakes training data pipelines."),
]

for i, (title, body) in enumerate(ml_items):
    col = i % 2
    row = i // 2
    x = 0.30 + col * 6.55
    y = 2.42 + row * 1.60
    is_crypto = "Cryptographic" in title
    is_combo  = "Integration" in title
    bg = GREEN_CHK if is_crypto else (TEAL if is_combo else RGBColor(0x22, 0x45, 0x70))
    add_rect(s, x, y, 6.30, 1.48, fill=bg)
    add_text(s, title, x+0.15, y+0.06, 6.05, 0.34, size=12, bold=True, color=GOLD)
    add_text(s, body,  x+0.15, y+0.42, 6.05, 0.98, size=10, color=LIGHT_BLUE)

# ── SLIDE 10: Cybersecurity & Forensics ───────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Cybersecurity & Forensics",
           "Post-hoc attribution, runtime protection, and the limits of non-cryptographic approaches")
footer(s)

forensic_items = [
    ("NTFS Journals\n($UsnJrnl, $LogFile)",
     "Update Sequence Number journal records every filesystem write at the volume level.\n"
     "$LogFile is the NTFS transaction log for crash recovery.\n"
     "Used in forensic investigation to reconstruct which files were modified and by which process.\n"
     "Limitation: provides attribution evidence post-hoc; does not prevent tampering\n"
     "and requires trusting the host OS and storage administrator."),
    ("Linux IMA / EVM",
     "Integrity Measurement Architecture: records SHA-1/SHA-256 of every executed file.\n"
     "Extended Verification Module: HMAC or signature over file metadata + content.\n"
     "Kernel-enforced: files not matching the stored hash are blocked from execution.\n"
     "Limitation: tied to a specific kernel instance; no portability across platforms."),
    ("Tamper-Tolerant Software (TTS)",
     "Combines obfuscation, self-checking guard threads, and graceful-degradation modes.\n"
     "Guards periodically compare pre-computed values against runtime state.\n"
     "On mismatch: crash, degrade, or alert rather than halt.\n"
     "Used in DRM, embedded systems, and anti-cheat software.\n"
     "Limitation: no provably secure software anti-tampering exists (Springer WISTP 2009)."),
    ("The Impossibility Result",
     "There are NO provably secure software anti-tampering methods.\n"
     "A proof of tamper-resistance requires:\n"
     "  1. A formal model capturing ALL attack surfaces (impossible for physical systems)\n"
     "  2. A reduction to an unsolved hard problem\n"
     "  3. That reduction to hold even when the attacker controls the hardware\n"
     "→ Cryptographic authentication makes a bounded claim within a stated adversary model."),
]

colors = [RGBColor(0x22, 0x45, 0x70), RGBColor(0x22, 0x45, 0x70),
          RGBColor(0x22, 0x45, 0x70), RGBColor(0x7A, 0x20, 0x20)]
ys = [1.38, 1.38, 1.38 + 2.58, 1.38 + 2.58]
xs = [0.30, 6.72, 0.30, 6.72]
ws = [6.25, 6.25, 6.25, 6.25]
hs = [2.45, 2.45, 2.45, 2.45]

for i, (title, body) in enumerate(forensic_items):
    add_rect(s, xs[i], ys[i], ws[i], hs[i], fill=colors[i])
    add_text(s, title, xs[i]+0.15, ys[i]+0.08, ws[i]-0.25, 0.40,
             size=13, bold=True, color=GOLD)
    add_text(s, body,  xs[i]+0.15, ys[i]+0.50, ws[i]-0.25, hs[i]-0.58,
             size=10, color=LIGHT_BLUE)

# ── SLIDE 11: Long-term Archival & Post-Quantum ────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Long-Term Archival & Post-Quantum Signatures",
           "Scientific archives signed today will be read — and attacked — decades from now")
footer(s)

add_rect(s, 0.35, 1.38, 12.63, 0.72, fill=RGBColor(0x7A, 0x40, 0x00))
add_text(s,
         "Shor's algorithm breaks RSA, ECDSA, Ed25519, and BLS pairings. "
         "A cryptographically relevant quantum computer in 2040–2050 can forge\n"
         "signatures over data harvested today — the 'harvest-now, forge-later' threat (T5).",
         0.55, 1.44, 12.2, 0.60, size=13, italic=True, color=RGBColor(0xFF, 0xE0, 0xA0))

pq_items = [
    ("NIST PQ Standards\n(FIPS 204 / 205)",
     "ML-DSA (Module-Lattice DSA, formerly Dilithium): lattice-based, 2420–4595 B signatures.\n"
     "SLH-DSA (formerly SPHINCS+): hash-based, stateless, 7856–49856 B signatures.\n"
     "Both standardized in 2024. Safe to deploy today for long-lived archives."),
    ("Hybrid Classical + PQ",
     "Ed25519 + ML-DSA-65: classical verifiability today + PQ security tomorrow.\n"
     "One hybrid signature per file (not per chunk) — storage overhead is negligible.\n"
     "Recommended transitional design per NIST and IETF drafts."),
    ("RFC 3161 Trusted Timestamps",
     "TSA (Timestamp Authority) countersigns a hash, proving it existed before a given time.\n"
     "Prevents rollback attacks (T4): verifier can reject stale versions using TSA evidence.\n"
     "Necessary for first-time recipients who have no prior version to compare against."),
    ("Sigstore Rekor (freshness log)",
     "Transparency log entry for each signed root: append-only, publicly auditable.\n"
     "Provides a tamper-evident temporal record of when a dataset was signed.\n"
     "Stretch goal in S2-D2 plan (§5.9): enables rollback detection for archived HDF5 files."),
    ("Why Merkle enables PQ upgrade",
     "The Merkle root is just a hash — swapping SHA-256 for SHA-3 or BLAKE3 is non-breaking.\n"
     "The signature over the root is the only PQ-sensitive component.\n"
     "Re-sign the root with ML-DSA without re-reading or re-hashing the entire dataset."),
    ("50-year archival horizon",
     "IllustrisTNG cosmological simulations: 75 TB, retained for decades.\n"
     "NOAA climate records: petabytes, legally required retention periods of 30–75 years.\n"
     "A signature placed today must be verifiable in 2075 — classical-only signatures are insufficient."),
]

for i, (title, body) in enumerate(pq_items):
    col = i % 2
    row = i // 2
    x = 0.30 + col * 6.55
    y = 2.24 + row * 1.68
    is_merkle = "Merkle" in title
    bg = GREEN_CHK if is_merkle else RGBColor(0x22, 0x45, 0x70)
    add_rect(s, x, y, 6.30, 1.56, fill=bg)
    add_text(s, title, x+0.15, y+0.06, 6.05, 0.42, size=12, bold=True, color=GOLD)
    add_text(s, body,  x+0.15, y+0.50, 6.05, 0.98, size=10, color=LIGHT_BLUE)

# ── SLIDE 12: Comparison across fields ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_GRAY)
header_bar(s, "Comparison: How Each Field Addresses Tampering",
           "Scientific data (current state) is uniquely underserved among fields that handle high-value, long-lived data")
footer(s)

comp_cols  = ["Field", "Primary technology", "Adversarial?", "Public verify?", "Subset proof?", "PQ-safe?", "Portable?"]
comp_col_x = [0.20, 2.15, 5.40, 6.70, 7.90, 9.10, 10.40]
comp_col_w = [1.85, 3.15, 1.20, 1.10, 1.10, 1.20, 2.70]

header_y = 1.28
for j, (lbl, x, w) in enumerate(zip(comp_cols, comp_col_x, comp_col_w)):
    add_rect(s, x, header_y, w-0.04, 0.40, fill=MID_BLUE)
    add_text(s, lbl, x+0.04, header_y+0.03, w-0.08, 0.34,
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

CHK, CRS, NA, PART = "✓", "✗", "—", "~"
comp_rows = [
    ("Storage / HPC",        "ZFS Merkle, dm-verity, T10-PI", CHK,  PART, CHK,  CRS,  CRS),
    ("Software supply chain","Code signing, Sigstore, CT logs",CHK,  CHK,  CRS,  PART, CHK),
    ("Blockchain",           "Merkle DAG, consensus ledger",  CHK,  CHK,  CHK,  CRS,  CHK),
    ("Scientific (today)",   "CRC, flat SHA-256",             CRS,  CRS,  CRS,  CRS,  PART),
    ("ML / AI",              "Anomaly detection, QML",        PART, NA,   NA,   NA,   PART),
    ("Forensics",            "NTFS journals, IMA/EVM",        PART, CRS,  CRS,  NA,   CRS),
    ("Long-term archival",   "ML-DSA, SLH-DSA, RFC 3161",    CHK,  CHK,  NA,   CHK,  CHK),
    ("Scientific (S2-D2)",   "Merkle + hybrid PQ sig",        CHK,  CHK,  CHK,  CHK,  CHK),
]

ROW_H = 0.60
for i, row in enumerate(comp_rows):
    y = header_y + 0.42 + i * ROW_H
    is_today    = "today" in row[0]
    is_proposed = "S2-D2" in row[0]
    bg = (RGBColor(0xFF, 0xE8, 0xE8) if is_today else
          RGBColor(0xC8, 0xE8, 0xC8) if is_proposed else
          (RGBColor(0xEB, 0xF2, 0xFB) if i % 2 == 0 else WHITE))
    add_rect(s, comp_col_x[0], y,
             sum(comp_col_w) + len(comp_col_w)*0.04, ROW_H-0.04, fill=bg)
    for j, (cell, x, w) in enumerate(zip(row, comp_col_x, comp_col_w)):
        c = DARK_BLUE
        if is_today and j == 0:    c = RED_CROSS
        if is_proposed and j == 0: c = GREEN_CHK
        if cell == CHK:  c = GREEN_CHK
        elif cell == CRS: c = RED_CROSS
        elif cell == PART: c = AMBER
        add_text(s, cell, x+0.04, y+0.10, w-0.08, ROW_H-0.18,
                 size=9, bold=(j == 0),
                 color=c,
                 align=PP_ALIGN.CENTER if j > 1 else PP_ALIGN.LEFT)

add_text(s, "~ = partial or probabilistic  |  — = not applicable  |  PART (Public) = vendor-specific or infrastructure-bound",
         0.20, 7.01, 12.9, 0.22, size=7.5, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 13: Conclusion ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Conclusion: Convergent Answer Across Fields",
           "Every field that needs portable, long-lived, publicly verifiable integrity converges on the same primitive")
footer(s)

add_rect(s, 0.35, 1.38, 12.63, 0.72, fill=GREEN_CHK)
add_text(s,
         "✓   Merkle tree  +  hybrid classical / post-quantum signature  =  the convergent solution",
         0.55, 1.50, 12.2, 0.54, size=18, bold=True, color=WHITE)

conclusions = [
    ("Every domain already uses Merkle trees",
     "Git, Bitcoin, ZFS, IPFS, npm, Certificate Transparency, and Android Verified Boot all rest\n"
     "on Merkle trees. S2-D2 applies the same well-proven primitive to scientific HDF5 data."),
    ("The gap is at the intersection of requirements",
     "No single field has needed: sub-dataset partial verification + public verifiability +\n"
     "post-quantum safety + no trusted setup + HPC-feasible build — simultaneously.\n"
     "Scientific long-lived archives are the first setting that demands all five."),
    ("ML detection and Merkle auth are complementary",
     "ML anomaly detectors (75–95 % accuracy) catch semantically suspicious changes;\n"
     "Merkle authentication detects ANY byte change with cryptographic certainty.\n"
     "Neither alone is sufficient for high-stakes ML training corpora."),
    ("The impossibility result focuses the claim",
     "There are no provably secure software anti-tampering methods.\n"
     "Merkle provides a cryptographic commitment within a stated adversary model — not\n"
     "a physical tamper-proof seal. That bounded claim is the right one for a software system."),
    ("Post-quantum is non-negotiable for 50-year archives",
     "Classical-only signatures on data archived today are at harvest-now, forge-later risk.\n"
     "NIST FIPS 204 (ML-DSA) + FIPS 205 (SLH-DSA) are standardized and deployable now.\n"
     "Merkle makes PQ upgrade cheap: re-sign only the root, not the entire dataset."),
]

for i, (title, body) in enumerate(conclusions):
    y = 2.28 + i * 1.02
    add_rect(s, 0.35, y, 0.06, 0.88, fill=GOLD)
    add_text(s, title, 0.52, y+0.02, 12.0, 0.32, size=13, bold=True, color=WHITE)
    add_text(s, body,  0.52, y+0.36, 12.0, 0.58, size=11, color=LIGHT_BLUE)

# ── save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "Tamper-Detection-Technologies.pptx")
prs.save(out_path)
print("Saved:", out_path)
