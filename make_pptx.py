"""
Build Sections 3 & 4 slide deck: Survey, Requirements, and Selection
Run with: /tmp/pptx-venv/bin/python3 make_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ── palette ──────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x37, 0x5E)
MID_BLUE   = RGBColor(0x27, 0x5D, 0x9B)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF7)
GOLD       = RGBColor(0xF0, 0xA5, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
RED_CROSS  = RGBColor(0xC0, 0x20, 0x20)
GREEN_CHK  = RGBColor(0x1A, 0x7A, 0x3C)
AMBER      = RGBColor(0xB8, 0x6A, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

# ── helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def slide_bg(slide, color=DARK_BLUE):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.25, fill=MID_BLUE)
    add_text(slide, title, 0.35, 0.12, 12.5, 0.65,
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.72, 12.5, 0.45,
                 size=15, color=LIGHT_BLUE)

def footer(slide, txt="S2-D2 · Tamper Detection & Data Integrity · Sections 3 & 4"):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill=MID_BLUE)
    add_text(slide, txt, 0.2, 7.17, 12.9, 0.28,
             size=9, color=WHITE, align=PP_ALIGN.CENTER)

def bullet_box(slide, items, l, t, w, h,
               size=15, color=WHITE, indent="  •  "):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = indent + item
        run.font.size = Pt(size)
        run.font.color.rgb = color

# ── SLIDE 1: Title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
add_rect(s, 0, 2.6, 13.33, 0.06, fill=GOLD)
add_text(s, "Tamper Detection and Data Integrity",
         0.6, 1.1, 12.0, 0.8, size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "for Self-Describing Scientific Data",
         0.6, 1.75, 12.0, 0.7, size=20, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "Sections 3 & 4 — Survey, Requirements, and Selection",
         0.6, 2.85, 12.0, 0.85, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Which integrity primitive survives 50-year scientific archives?",
         0.6, 3.75, 12.0, 0.6, size=18, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "M. Scot Breitenfeld · The HDF Group · S2-D2 Research Plan",
         0.6, 6.45, 12.0, 0.5, size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 2: Agenda ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Roadmap for This Section")
footer(s)

steps = [
    ("§3", "Survey of integrity primitives",    "14 families — capability and killer flaw for HDF5"),
    ("§3", "What the survey reveals",           "Four structural barriers that drive the requirements"),
    ("§4", "Threat model (T1–T8)",              "Eight threats that motivate the requirements"),
    ("§4", "Twelve requirements (R1–R12)",      "What any primitive must satisfy for HDF5"),
    ("§4", "Comparison table",                  "All twelve requirements scored across seven alternatives"),
    ("§4", "Selection: Merkle + PQ signature",  "The only R1–R12-complete, quantum-safe option"),
]
for i, (num, title, sub) in enumerate(steps):
    y = 1.45 + i * 0.92
    add_rect(s, 0.4, y, 0.65, 0.55, fill=GOLD)
    add_text(s, num, 0.40, y+0.06, 0.65, 0.48,
             size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_rect(s, 1.18, y, 11.5, 0.55, fill=MID_BLUE)
    add_text(s, title, 1.28, y+0.04, 5.8, 0.48, size=17, bold=True, color=WHITE)
    add_text(s, sub,   7.2,  y+0.07, 5.3, 0.42, size=13, color=LIGHT_BLUE)

# ── SLIDE 3: Survey landscape ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_GRAY)
header_bar(s, "§3  Survey: 14 Integrity Primitive Families",
           "Each surveyed for suitability on multi-terabyte, long-lived, chunked scientific data")
footer(s)

# Two-column table: family | examples | killer flaw
families = [
    ("Error-detection codes",       "CRC-32, Fletcher-32, Adler-32, T10-PI",
     "Invertible in closed form — cannot resist adversarial tampering"),
    ("Cryptographic hashes",        "SHA-256, BLAKE3, SHA-3, KangarooTwelve",
     "Flat hash: verifying any chunk requires re-reading every byte"),
    ("MACs / AEAD",                 "HMAC-SHA-256, ChaCha20-Poly1305, AES-GCM",
     "Secret key required — third-party auditors and collaborators locked out"),
    ("Classical signatures",        "Ed25519, ECDSA P-256, RSA-PSS",
     "Broken by Shor's algorithm — 50-year archives signed today are at risk"),
    ("Post-quantum signatures",     "ML-DSA-65, SLH-DSA, XMSS",
     "Per-chunk: 33 GB of sigs + 3 CPU-hours to sign 10⁷ chunks"),
    ("Hash trees / ADS",            "Merkle binary, SMT, Merkle Mountain Range",
     "→ THE CANDIDATE  (the only family that meets all §4 requirements)"),
    ("Verkle / KZG commitments",    "KZG polynomial, Ethereum Verkle tree",
     "Trusted-setup ceremony required; pairing broken by Shor"),
    ("Cryptographic accumulators",  "RSA accumulator, bilinear accumulator",
     "Trusted setup + not PQ + set-based (no chunk order)"),
    ("ZK-SNARKs / Folding",         "Groth16, PLONK, Nova, HyperNova",
     "O(N) proving cost; > 24 h to commit 10⁷ chunks"),
    ("Transparency logs",           "Certificate Transparency, Sigstore Rekor",
     "Append-only — HDF5 workloads overwrite chunks in place"),
    ("TEE / fs-verity",             "Intel SGX, ARM TrustZone, Linux fs-verity",
     "Vendor-locked; stripped when file is transferred or archived"),
    ("FS / cloud checksums",        "ZFS, Btrfs, S3 ETag, GCS CRC32C",
     "Stripped in transfer; no partial verification; no non-repudiation"),
    ("Scientific format checksums", "NetCDF CDF-5, Zarr v3 chunk hashes, Parquet",
     "Non-adversarial; no public signing or subset-proof capability"),
    ("W3C PROV / lineage",          "PROV-DM, ProvONE, PROV-JSON",
     "Provenance graph — complements integrity but is not a substitute"),
]

# Split into two columns of 7
left  = families[:7]
right = families[7:]
ROW_H = 0.76

for col_idx, group in enumerate([left, right]):
    lx = 0.20 if col_idx == 0 else 6.78
    for i, (name, examples, flaw) in enumerate(group):
        y = 1.38 + i * ROW_H
        is_candidate = "CANDIDATE" in flaw
        bar_color = GREEN_CHK if is_candidate else MID_BLUE
        add_rect(s, lx, y, 6.38, ROW_H - 0.06, fill=bar_color)
        add_text(s, name,    lx+0.12, y+0.03,  6.1, 0.26,
                 size=12, bold=True,
                 color=GOLD if is_candidate else WHITE)
        add_text(s, examples, lx+0.12, y+0.27, 6.1, 0.20,
                 size=9, italic=True,
                 color=LIGHT_BLUE if not is_candidate else RGBColor(0xC0,0xFF,0xC0))
        add_text(s, flaw,    lx+0.12, y+0.46,  6.1, 0.24,
                 size=9,
                 color=GOLD if is_candidate else RGBColor(0xFF, 0xCC, 0xAA))

# ── SLIDE 4: What the survey reveals ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "§3  What the Survey Reveals",
           "Four structural barriers that no alternative escapes simultaneously")
footer(s)

barriers = [
    (
        "1",
        "Non-cryptographic codes are invertible",
        "CRC, Fletcher, and Adler all admit closed-form corrections.\n"
        "An adversary computes a matching tail that preserves the checksum — no brute force needed.\n"
        "→  The integrity primitive must be cryptographic at its core.",
    ),
    (
        "2",
        "Symmetric auth requires a shared secret — rules out public verifiability",
        "MACs and AEAD are ideal per-chunk tools, but the verifier must hold the key.\n"
        "Third-party auditors, archive mirrors, and collaborators cannot verify.\n"
        "→  Public verifiability eliminates any purely symmetric scheme as the root.",
    ),
    (
        "3",
        "Pairing- and factoring-based schemes fail the 50-year archival horizon",
        "Shor's algorithm breaks ECDSA, Ed25519, RSA, BLS12-381 pairings, and KZG commitments.\n"
        "Scientific archives are signed today and verified in 2040–2070.\n"
        "→  The post-quantum requirement eliminates Verkle, BLS aggregation, and RSA/bilinear accumulators.",
    ),
    (
        "4",
        "Trusted-setup ceremonies are unacceptable for long-lived archives",
        "KZG, Verkle, and RSA accumulators require an MPC ceremony whose\n"
        "parameters must be trusted for the archive's entire lifetime.\n"
        "→  No-trusted-setup eliminates every scheme with a structured reference string.",
    ),
]

for i, (num, title, body) in enumerate(barriers):
    y = 1.42 + i * 1.42
    add_rect(s, 0.35, y, 0.62, 1.25, fill=GOLD)
    add_text(s, num, 0.35, y+0.35, 0.62, 0.55,
             size=26, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_rect(s, 1.08, y, 11.9, 1.25, fill=MID_BLUE)
    add_text(s, title, 1.20, y+0.06, 11.6, 0.36,
             size=14, bold=True, color=WHITE)
    add_text(s, body,  1.20, y+0.44, 11.6, 0.74,
             size=11, color=LIGHT_BLUE)

# ── SLIDE 5: Threat model ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "§4  Threat Model: Eight Threats That Drive the Requirements",
           "Adversary model defined before requirements — threats first, then what the primitive must resist")
footer(s)

threats = [
    ("T1", "Storage-level tampering",   "TCD",
     "Adversary with storage access modifies chunk data or metadata — targeted or wholesale replacement"),
    ("T2", "Silent data corruption",    "FMT",
     "Hardware bit-flips, firmware bugs, or filesystem errors corrupt chunks undetected"),
    ("T3", "Provenance forgery",        "EXT",
     "Adversary attributes a modified dataset to a different creator — requires non-repudiable signatures"),
    ("T4", "Rollback / stale data",     "sys",
     "Valid, properly signed older version replaces current file — cryptography passes but data is stale"),
    ("T5", "Harvest-now, forge-later",  "PQ",
     "Archives signed today are broken by CRQC in 2040–2070 via Shor's algorithm on Ed25519 / ECDSA"),
    ("T6", "Algorithm downgrade",       "sys",
     "Adversary strips provenance metadata or degrades hash/AEAD identifiers to force fail-open or weaker crypto"),
    ("T7", "Verification DoS",          "sys",
     "Malformed companion dataset with crafted tree depth or coverage certificate exhausts verifier memory"),
    ("T8", "Structural leakage",        "priv",
     "SMT null-leaf positions reveal unallocated chunk grid — geographic or genomic structure exposed without decryption"),
]

left  = threats[:4]
right = threats[4:]
ROW_H_T = 1.28

CAT_COLORS = {
    "TCD":  RGBColor(0xB0, 0x30, 0x20),
    "FMT":  RGBColor(0xB0, 0x30, 0x20),
    "EXT":  RGBColor(0xB0, 0x30, 0x20),
    "sys":  RGBColor(0x18, 0x55, 0x8A),
    "PQ":   RGBColor(0x7A, 0x40, 0x00),
    "priv": RGBColor(0x2A, 0x6A, 0x2A),
}

for col_idx, group in enumerate([left, right]):
    lx = 0.25 if col_idx == 0 else 6.72
    for i, (tag, name, cat, desc) in enumerate(group):
        y = 1.40 + i * ROW_H_T
        cat_color = CAT_COLORS.get(cat, MID_BLUE)
        add_rect(s, lx, y, 6.30, ROW_H_T - 0.08, fill=MID_BLUE)
        add_rect(s, lx, y, 0.60, ROW_H_T - 0.08, fill=GOLD)
        add_text(s, tag, lx, y + 0.34, 0.60, 0.50,
                 size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        add_rect(s, lx + 4.70, y + 0.08, 1.48, 0.36, fill=cat_color)
        add_text(s, cat, lx + 4.72, y + 0.09, 1.44, 0.32,
                 size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, name, lx + 0.68, y + 0.06, 4.1, 0.36,
                 size=13, bold=True, color=WHITE)
        add_text(s, desc, lx + 0.68, y + 0.52, 5.50, 0.60,
                 size=10, color=LIGHT_BLUE)

# ── SLIDE 6: Requirements overview ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "§4  Twelve Requirements for HDF5 Integrity",
           "Any candidate primitive must satisfy all twelve — simultaneously")
footer(s)

reqs = [
    ("R1",  "Partial verification",        "Verify k chunks at O(k log N), no full rehash"),
    ("R2",  "Incremental extend",          "Append chunks in O(log N), no re-sign"),
    ("R3",  "In-place update",             "Overwrite chunk k in O(log N) tree-update work"),
    ("R4",  "Tamper localization",         "Identify which chunks changed, not just 'something did'"),
    ("R5a", "Chunk authenticity",          "Each chunk tied to signer's commitment"),
    ("R5b", "Coverage proof",              "Prove delivered set is exactly the requested hyperslab"),
    ("R6",  "Public verifiability",        "No secret key required — any auditor can verify"),
    ("R7",  "Compact proof",               "Proof size sublinear in N  (target: O(k log N))"),
    ("R8",  "No trusted setup",            "No MPC ceremony, no trapdoored parameters"),
    ("R9",  "Post-quantum path",           "Survives Shor: swap hash/sig, no re-sign of archives"),
    ("R10", "HPC-feasible build",          "Commitment at memory-bandwidth speed on commodity HW"),
    ("R11", "Distributed parallel build",  "Leaves computed independently per node; O(log P) MPI reduction"),
    ("R12", "Structure confidentiality",   "Sparse null-leaf positions hidden from observers who cannot decrypt"),
]
cols = [(0.35, reqs[:6]), (6.85, reqs[6:])]
for cx, group in cols:
    for i, (tag, name, detail) in enumerate(group):
        y = 1.42 + i * 0.90
        add_rect(s, cx, y, 0.72, 0.56, fill=GOLD)
        add_text(s, tag, cx, y+0.06, 0.72, 0.46,
                 size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        add_rect(s, cx+0.76, y, 5.7, 0.56, fill=MID_BLUE)
        add_text(s, name,   cx+0.86, y+0.02, 5.5, 0.26, size=12, bold=True, color=WHITE)
        add_text(s, detail, cx+0.86, y+0.28, 5.5, 0.24, size=10, color=LIGHT_BLUE)

# ── SLIDE 7: Comparison table ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_GRAY)
header_bar(s, "§4  Candidate Primitives vs. Requirements")
footer(s)

headers = ["Req.", "Flat\nhash", "Per-chk\nMAC", "Per-chk\nsig.", "RSA\nacc.",
           "KZG/\nVerkle", "ZKP/\nPoR", "TEE/\nfs-verity", "Merkle\ntree ✓"]
col_w   = [1.05, 1.12, 1.12, 1.12, 1.12, 1.12, 1.12, 1.20, 1.20]
col_x   = [0.18]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

ROW_H = 0.350
header_y = 1.32

for j, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
    fill = DARK_BLUE if j == len(headers)-1 else MID_BLUE
    add_rect(s, x, header_y, w-0.04, 0.56, fill=fill)
    add_text(s, hdr, x+0.03, header_y+0.02, w-0.07, 0.52,
             size=10, bold=True, color=WHITE if j != len(headers)-1 else GOLD,
             align=PP_ALIGN.CENTER)

CHK = "✓"
CRS = "✗"
NA  = "—"

rows = [
    ("R1",  CRS, CHK, CHK, CHK, CHK, "samp", CHK,   CHK),
    ("R2",  CHK, CHK, CHK, CRS, CRS, CRS,   CHK,   CHK),
    ("R3",  CRS, CHK, CHK, CRS, CRS, CRS,   CHK,   CHK),
    ("R4",  CRS, CHK, CHK, CRS, "prt", CRS, CHK,   CHK),
    ("R5a", CRS, CHK, CHK, CHK, CHK, CHK,   CHK,   CHK),
    ("R5b", CRS, CRS, CRS, CRS, CHK, CHK,   CRS,   CHK),
    ("R6",  CHK, CRS, CHK, CHK, CHK, CHK,   "vnd", CHK),
    ("R7",  NA,  "16k","64k","O(1)","O(1)","O(1)", NA,  "32k·lgN"),
    ("R8",  CHK, CHK, CHK, CRS, CRS, "var", "vnd", CHK),
    ("R9",  CHK, CHK, CRS, CRS, CRS, "STARK","vnd", CHK),
    ("R10", CHK, CHK, CRS, CRS, "1s", CRS,  CHK,   CHK),
    ("R11", CHK, CHK, CHK, CRS, CRS, CRS,  "vnd", CHK),
    ("R12", CRS, CRS, CRS, CRS, CRS, CRS,  CRS,   "opt"),
]

for i, row in enumerate(rows):
    y = header_y + 0.58 + i * ROW_H
    bg = RGBColor(0xEB, 0xF2, 0xFB) if i % 2 == 0 else WHITE
    add_rect(s, col_x[0], y, sum(col_w)-0.04, ROW_H-0.02, fill=bg)
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        if j == len(row)-1:
            add_rect(s, x, y, w-0.04, ROW_H-0.02, fill=RGBColor(0xD0, 0xEA, 0xD0))
        txt_color = DARK_BLUE
        if cell == CHK:   txt_color = GREEN_CHK
        elif cell == CRS: txt_color = RED_CROSS
        add_text(s, cell, x+0.02, y+0.04, w-0.06, ROW_H-0.08,
                 size=11, bold=(cell in (CHK, CRS)),
                 color=txt_color, align=PP_ALIGN.CENTER)

add_text(s, "samp=sampling only  |  prt=partial  |  vnd=vendor-locked  |  var=varies  |  opt=opt-in (SMT with PRF-masked null hashes)",
         0.18, 7.00, 12.9, 0.28, size=8.5, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 8: Selection conclusion ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "§4  Selection: Hash Tree with a Signed Root",
           "The only primitive satisfying R1–R12 without trusted setup or quantum liability")
footer(s)

add_rect(s, 0.35, 1.38, 12.63, 0.72, fill=GREEN_CHK)
add_text(s,
         "✓   Merkle tree (binary / sparse)  +  hybrid Ed25519 · ML-DSA-65 root signature",
         0.55, 1.44, 12.2, 0.62, size=19, bold=True, color=WHITE)

points = [
    ("Subset & coverage proofs (R5a, R5b, R7)",
     "O(k log N) multiproof — comes directly from tree structure; no extra primitive needed"),
    ("Tamper localization (R4)",
     "Root-to-leaf walk isolates modified chunks in O(log N)"),
    ("In-place update & append (R2, R3)",
     "Recompute only ⌈log₂ N⌉ ≈ 24 SHA-256 calls per update at N = 10⁷  (~3 µs on SHA-NI)"),
    ("No trusted setup · Post-quantum (R8, R9)",
     "PQ-safe hash today; swap signature to ML-DSA-65 / SLH-DSA without re-signing archives"),
    ("HPC-feasible build (R10)",
     "10 TB committed in ~30 s across ~200 HPC cores; tree aggregation < 1 s on one core"),
    ("Public verifiability (R6)",
     "Only the signer's public key needed — any auditor, mirror, or collaborator can verify"),
]
for i, (title, detail) in enumerate(points):
    y = 2.24 + i * 0.79
    add_rect(s, 0.35, y, 0.06, 0.58, fill=GOLD)
    add_text(s, title,  0.52, y+0.02, 12.0, 0.30, size=14, bold=True, color=WHITE)
    add_text(s, detail, 0.52, y+0.31, 12.0, 0.30, size=12, color=LIGHT_BLUE)

# ── SLIDE 9: Why not the alternatives ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "§4  Why Not the Alternatives?",
           "Each alternative falls to a single knockout criterion")
footer(s)

alts = [
    ("Verkle trees\n(KZG-based)",
     "R8 + R9",
     "KZG requires a trusted-setup ceremony. The underlying elliptic-curve pairing is broken\n"
     "by Shor's algorithm. Re-signing 30 years of NOAA records in 2045 is not an option."),
    ("BLS-aggregated\nper-chunk sigs",
     "R9 + R5b",
     "BLS12-381 pairings are broken by Shor. No upgrade path exists.\n"
     "Coverage proof (R5b) still needs a signed manifest — which reduces to a Merkle root."),
    ("Per-chunk PQ sigs\n(ML-DSA-65)",
     "R7 + R10",
     "33 GB of signatures per dataset (N = 10⁷ chunks). Signing takes ~3 CPU-hours.\n"
     "Two orders of magnitude worse than Merkle with no compensating benefit."),
    ("ZK-SNARKs /\nFolding schemes",
     "R10",
     "O(N) field-operation proving overhead for static data subsetting.\n"
     "Designed for computation correctness, not data authentication. > 24 h to commit."),
    ("RSA accumulators",
     "R8 + R9",
     "Require trusted setup (trapdoor exponent). Broken by Shor.\n"
     "No advantage over Merkle on any other dimension."),
    ("Per-chunk MACs\n/ AEAD only",
     "R6",
     "Verification requires the secret key. Third-party auditors and\n"
     "downstream researchers are locked out by design."),
]
for i, (name, fail, reason) in enumerate(alts):
    col = i % 2
    row = i // 2
    x = 0.35 + col * 6.55
    y = 1.42 + row * 1.87
    add_rect(s, x, y, 6.2, 1.72, fill=RGBColor(0x22, 0x45, 0x70))
    add_rect(s, x+4.55, y+0.10, 1.50, 0.42, fill=RED_CROSS)
    add_text(s, "Fails " + fail, x+4.58, y+0.12, 1.44, 0.38,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, name,   x+0.15, y+0.08,  4.3, 0.55, size=14, bold=True, color=GOLD)
    add_text(s, reason, x+0.15, y+0.62, 5.95, 0.98, size=11, color=LIGHT_BLUE)

# ── SLIDE 10: Honest framing / conclusion ────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "§4  Honest Framing of the Conclusion")
footer(s)

add_rect(s, 0.35, 1.38, 12.63, 1.10, fill=RGBColor(0x18, 0x45, 0x30))
add_text(s,
         "Merkle is not the only primitive satisfying R1–R7 in absolute terms.\n"
         "Verkle and BLS-aggregated signatures are competitive there.",
         0.55, 1.44, 12.2, 1.0, size=15, italic=True, color=RGBColor(0xA8, 0xD8, 0xB8))

add_rect(s, 0.35, 2.65, 12.63, 1.28, fill=MID_BLUE)
add_text(s, "What distinguishes Merkle:", 0.55, 2.70, 12.0, 0.38,
         size=16, bold=True, color=WHITE)
add_text(s,
         "It is the only primitive achieving R1–R12 simultaneously,\n"
         "without a trusted-setup ceremony and without a pairing- or\n"
         "factoring-based primitive that Shor's algorithm breaks.",
         0.55, 3.06, 12.0, 0.80, size=15, color=WHITE)

add_rect(s, 0.35, 4.10, 12.63, 0.06, fill=GOLD)

add_text(s, "R8 and R9 are the discriminating requirements — forced by the multi-decade\n"
            "archival horizon of scientific data. This constraint is absent from the\n"
            "blockchain settings where KZG, Verkle, and BLS were designed.",
         0.55, 4.25, 12.2, 0.90, size=14, italic=True, color=LIGHT_BLUE)

add_text(s, "Design choices that follow from this selection:",
         0.55, 5.30, 12.0, 0.38, size=14, bold=True, color=WHITE)
choices = [
    "AEAD inside the Encrypt-then-MAC filter pipeline  (§5.7)",
    "Classical + PQ hybrid root signature: Ed25519 + ML-DSA-65  (§5.8)",
    "Transparency-log publication of signed roots for freshness  (§5.9)",
]
bullet_box(s, choices, 0.55, 5.65, 12.0, 1.20, size=13, color=LIGHT_BLUE)

# ── save ──────────────────────────────────────────────────────────────────────
out = "/home/brtnfld/work/R2D2-Yr2/Section4-Requirements-and-Selection.pptx"
prs.save(out)
print("Saved:", out)
